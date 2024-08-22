import json
from concurrent.futures import ThreadPoolExecutor, as_completed
from multiprocessing import cpu_count
from django.http import JsonResponse, HttpResponse
from django.shortcuts import render
from django.conf import settings
from django.views.decorators.csrf import csrf_exempt
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

import os
os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'HICC_PROJECT.settings')

from django.conf import settings
import openai

# OpenAI API 키 설정
openai.api_key = settings.OPENAI_API_KEY

# 현재 기기의 코어 개수를 기반으로 최대 병렬 처리 개수 설정
max_workers = cpu_count()


def home(request):
    return render(request, 'mainapp/home.html')


@csrf_exempt
def translate(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        source_language = data['source_language']
        target_language = data['target_language']
        text = data['text']
        print(f"Received text: {text}")  # 디버깅용 임시코드

        if not text:
            return JsonResponse({'translated_text': '텍스트가 비어 있습니다.'}, status=400)

        translated_text = translate_text(text, source_language, target_language)
        print(f"Translated text: {translated_text}")  # 디버깅용 임시코드
        return JsonResponse({'translated_text': translated_text})
    return render(request, 'mainapp/text_tr_improved.html')


def translate_text(text, source_language, target_language):
    if not text.strip():
        return text

    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo-1106",
            messages=[
                {"role": "system",
                 "content": f"You are a translation assistant that translates {source_language} text to {target_language}."},
                {"role": "user",
                 "content": f"Translate the following text from {source_language} to {target_language}: {text}"}
            ],
            max_tokens=1024,
            temperature=0.5,
        )
        print(f"API Response: {response}")  # 디버깅용 임시코드
        translated_text = response.choices[0].message['content'].strip()
        if not translated_text:
            return '번역 실패'
        return translated_text
    except Exception as e:
        print(f"Error: {e}")
        return '번역 실패'

@csrf_exempt
def summarize_text(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        text = data.get('text')
        if not text:
            return JsonResponse({'summary': '텍스트가 비어 있습니다.'}, status=400)

        summarized_text = get_summary(text)
        return JsonResponse({'summary': summarized_text})
    return JsonResponse({'error': '잘못된 요청입니다.'}, status=400)


def get_summary(text):
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a summary assistant that summarizes text."},
                {"role": "user", "content": f"Summarize the following text in Korean: {text}"}
            ],
            max_tokens=512,
            temperature=0.5,
        )
        summarized_text = response.choices[0].message['content'].strip()
        return summarized_text
    except Exception as e:
        print(f"Error: {e}")
        return '요약 실패'

@csrf_exempt
def define_term(request):
    term = request.GET.get('term')
    if not term:
        return JsonResponse({'definition': '용어가 제공되지 않았습니다.'}, status=400)

    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are an expert assistant that provides detailed definitions of technical terms in Korean."},
                {"role": "user", "content": f"Define the following term in Korean: {term}"}
            ],
            max_tokens=512,
            temperature=0.5,
        )
        definition = response.choices[0].message['content'].strip()
        return JsonResponse({'definition': definition})
    except Exception as e:
        print(f"Error: {e}")
        return JsonResponse({'definition': '정의 실패'}, status=500)


def translate_text_with_context(text, source_language, target_language):
    # 텍스트가 너무 짧거나 의미가 없을 경우 번역하지 않고 원본을 반환
    if len(text.strip()) <= 4:
        return text

    prompt = (
        f"Translate the following text from {source_language} to {target_language}. "
        "If the text contains technical terms, names, or specific jargon, do not translate those terms. "
        f"Here is the text: {text}"
    )

    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a professional translator."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=1024,
            temperature=0.5,
        )
        translated_text = response.choices[0].message['content'].strip()
        return translated_text
    except Exception as e:
        print(f"Error: {e}")
        return '번역 실패'

def adjust_text_in_box(shape, translated_text, slide):
    text_frame = shape.text_frame
    original_font_size = None
    original_alignment = None
    original_width = shape.width  # 원본 텍스트 박스 너비 저장
    original_height = shape.height  # 원본 텍스트 박스 높이 저장

    # 원본 텍스트 속성 보존
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        first_run = text_frame.paragraphs[0].runs[0]
        original_font_size = first_run.font.size  # 원본 텍스트 크기 저장
        original_alignment = text_frame.paragraphs[0].alignment  # 원본 정렬 저장

    # 기존 텍스트를 완전히 제거
    text_frame.clear()

    # 첫 번째 단락에 번역된 텍스트 추가
    p = text_frame.paragraphs[0]  # 기존 단락을 가져옴
    run = p.add_run()
    run.text = translated_text  # 번역된 텍스트 추가

    # 원본 글자 크기와 정렬 설정
    if original_font_size:
        run.font.size = original_font_size
    if original_alignment:
        p.alignment = original_alignment

    def is_overlapping(new_shape):
        """기존 다른 요소들과 겹치는지 확인"""
        for other_shape in slide.shapes:
            if other_shape == shape:
                continue
            if (new_shape.left < other_shape.left + other_shape.width and
                new_shape.left + new_shape.width > other_shape.left and
                new_shape.top < other_shape.top + other_shape.height and
                new_shape.top + new_shape.height > other_shape.top):
                return True
        return False

    # 텍스트 상자가 텍스트 프레임보다 작을 경우 크기 조정
    while True:
        # 텍스트 프레임이 상자에 잘 맞는지 확인
        text_box_width = shape.width
        text_box_height = shape.height

        # 텍스트가 상자에 맞으면 루프 종료
        if text_box_width >= shape.width and text_box_height >= shape.height:
            break

        # 텍스트 상자의 너비와 높이를 늘려서 텍스트가 들어갈 수 있게 조정
        shape.width += Inches(0.1)  # 너비를 약간 늘림
        shape.height += Inches(0.1)  # 높이를 약간 늘림

        # 겹치는지 확인
        if is_overlapping(shape):
            # 겹치면 크기를 원래대로 돌리고 루프 종료
            shape.width = original_width
            shape.height = original_height
            break

        # 텍스트 상자가 너무 커지면 제한
        if shape.width >= original_width * 2 or shape.height >= original_height * 2:
            break


def translate_shape_text(shape, slide, source_language, target_language):
    if shape.has_text_frame and shape.text_frame.text.strip():
        original_text = shape.text_frame.text.strip()
        if original_text:
            translated_text = translate_text_with_context(original_text, source_language, target_language)
            adjust_text_in_box(shape, translated_text, slide)
    return shape
@csrf_exempt
def pptTranslate(request):
    if request.method == 'POST' and request.FILES.get('ppt_file'):
        ppt_file = request.FILES['ppt_file']
        source_language = request.POST.get('source_language')
        target_language = request.POST.get('target_language')
        presentation = Presentation(ppt_file)
        translated_ppt = BytesIO()

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    # 이미지 등 텍스트가 아닌 요소는 제외
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        # translate_shape_text 함수에 slide 인자를 추가하여 전달
                        futures.append(executor.submit(translate_shape_text, shape, slide, source_language, target_language))

            for future in as_completed(futures):
                future.result()  # 각 작업의 결과를 가져옴 (여기서는 결과가 필요하지 않지만 예외 처리를 위해 사용)

        presentation.save(translated_ppt)
        translated_ppt.seek(0)

        response = HttpResponse(translated_ppt,
                                content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        response['Content-Disposition'] = 'attachment; filename="translated.pptx"'
        return response

    return render(request, 'mainapp/ppttranslate.html')







